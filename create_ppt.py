from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def create_presentation():
    prs = Presentation()
    
    # Colors
    BG_COLOR = RGBColor(15, 23, 42)      # Dark Navy
    TITLE_COLOR = RGBColor(255, 255, 255)# White
    ACCENT_BLUE = RGBColor(59, 130, 246) # Bright Blue
    ACCENT_ORANGE = RGBColor(245, 158, 11) # Amber
    TEXT_COLOR = RGBColor(226, 232, 240) # Light Grey

    def apply_design(slide, is_title_slide=False):
        # 1. Set Background to Dark Navy
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
        # 2. Add Decorative Bars
        if is_title_slide:
            # Center decorations
            left = Inches(3)
            top = Inches(0.5)
            width = Inches(4)
            height = Inches(0.1)
            shape = slide.shapes.add_shape(1, left, top, width, height) # Rectangle
            shape.fill.solid()
            shape.fill.fore_color.rgb = ACCENT_BLUE
            shape.line.fill.background()
        else:
            # Header Bar
            left = Inches(0.5)
            top = Inches(1.2)
            width = Inches(1.5)
            height = Inches(0.05)
            shape = slide.shapes.add_shape(1, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = ACCENT_ORANGE
            shape.line.fill.background()

    def set_text(shape, text, is_title=False, is_bold=False, font_size=None, color=TEXT_COLOR):
        shape.text = text
        for p in shape.text_frame.paragraphs:
            p.font.color.rgb = color
            p.font.name = 'Malgun Gothic' # Windows Korean Font
            if is_bold:
                p.font.bold = True
            if font_size:
                p.font.size = font_size

    # --- 1. Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    apply_design(slide, is_title_slide=True)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    set_text(title, "Office AI\nBusiness Plan", is_title=True, is_bold=True, font_size=Pt(54), color=TITLE_COLOR)
    set_text(subtitle, "KS 심사위원/정부지원 행정 업무 자동화를 위한\n중소기업 전용 AI 에이전트 서비스\n\n발표자: 노진광 대표", font_size=Pt(20))

    # --- Common Content Helper ---
    def add_content_slide(title_text, content_items):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        apply_design(slide)
        
        # Title
        title = slide.shapes.title
        set_text(title, title_text, is_title=True, is_bold=True, font_size=Pt(36), color=TITLE_COLOR)
        title.top = Inches(0.3)
        
        # Content
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear() # Clear default
        
        for item in content_items:
            p = tf.add_paragraph()
            p.text = item['text']
            p.font.name = 'Malgun Gothic'
            p.font.color.rgb = TEXT_COLOR
            p.level = item.get('level', 0)
            
            # Styling
            if item.get('bold'):
                p.font.bold = True
                p.font.color.rgb = ACCENT_BLUE
                p.font.size = Pt(24)
            else:
                p.font.size = Pt(20)
                
            # Spacing
            p.space_after = Pt(12)

    # --- 2. Problem ---
    add_content_slide("1. 문제 인식: 중소기업의 인력난", [
        {'text': "만성적인 전문 인력 부족 현상", 'bold': True},
        {'text': "• 국내 중소기업 80% 이상이 인력난 호소", 'level': 0},
        {'text': "• 기획/품질관리/인증 등 특수 직무 수행 인력 채용 불가", 'level': 0},
        {'text': ""},
        {'text': "핵심 Pain Point", 'bold': True},
        {'text': "• 복잡한 서류 작성 부담으로 정부지원사업(R&D 등) 포기", 'level': 0},
        {'text': "• 잦은 담당자 퇴사로 인한 KS 심사 준비 단절", 'level': 0}
    ])

    # --- 3. Solution ---
    add_content_slide("2. 해결 방안: Office AI", [
        {'text': "KS 심사위원이 직접 설계한 AI 비서", 'bold': True},
        {'text': "• 핵심 기능: 기업 데이터 입력 시 사업계획서/심사 서류 자동 생성", 'level': 0},
        {'text': "• 특화 기술: Domain-Specific RAG", 'level': 0},
        {'text': "  - KS 인증 심사 기준서 및 합격 데이터셋 학습", 'level': 1},
        {'text': "  - 심사위원 페르소나를 통한 'Self-Check' 기능", 'level': 1}
    ])

    # --- 4. Competitiveness ---
    add_content_slide("2. 해결 방안: 경쟁력", [
        {'text': "압도적인 도메인 지식 (Unfair Advantage)", 'bold': True},
        {'text': "• KSA 심사팀 출신 + 에너지 기업 CEO의 현장 경험 결합", 'level': 0},
        {'text': "• 단순 기술 기업은 모방할 수 없는 '심사 노하우' 보유", 'level': 0},
        {'text': ""},
        {'text': "사용자 편의성", 'bold': True},
        {'text': "• \"자금 필요해\" 같은 자연어 명령어로 워크플로우 실행", 'level': 0}
    ])

    # --- 5. BM ---
    add_content_slide("3. 비즈니스 모델 & 시장 진입", [
        {'text': "초기 타겟 시장 (Beachhead Market)", 'bold': True},
        {'text': "• 태양광 시공/제조 업체 (서류 작업 과다 업종)", 'level': 0},
        {'text': "• CEO의 4억 규모 계약 관리 네트워크 직접 활용", 'level': 0},
        {'text': ""},
        {'text': "수익 모델 (Revenue Stream)", 'bold': True},
        {'text': "• 월 구독료: 20만 원 (AI 행정 비서)", 'level': 0},
        {'text': "• 성공 보수: 정부 과제 선정 시 수수료 3~5%", 'level': 0}
    ])

    # --- 6. Team ---
    add_content_slide("4. 팀 역량 (Team)", [
        {'text': "대표자: 노진광 (Biz & Tech 융합 Leader)", 'bold': True},
        {'text': "• [Energy] 정도에코에너지/제이케이네트웍스 대표", 'level': 0},
        {'text': "  - 대규모 발전소 시공 및 운영 계약 총괄", 'level': 1},
        {'text': "• [Certification] KSA 인증 심사팀 출신", 'level': 0},
        {'text': "  - KS 인증/ISO 심사 프로세스 전문성", 'level': 1},
        {'text': "• [Digital] AI Server 구축 및 자동화 개발", 'level': 0},
        {'text': "  - NotebookLM, Claude 에이전트 서비스 자체 구현", 'level': 1}
    ])

    prs.save('C:\\work\\office-ai\\Office_AI_Business_Plan_v2.pptx')
    print("Succeed")

if __name__ == "__main__":
    create_presentation()
