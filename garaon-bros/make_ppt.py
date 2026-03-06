"""
가라온 브로즈 프로젝트 소개 PPT 생성 스크립트
실행: py make_ppt.py
출력: garaon-bros-presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── 색상 팔레트 ──────────────────────────────────
BG       = RGBColor(0x0b, 0x0b, 0x1a)   # 배경 (거의 검정)
CARD     = RGBColor(0x15, 0x15, 0x26)   # 카드 배경
CARD2    = RGBColor(0x1e, 0x1e, 0x38)   # 카드 배경 2
CYAN     = RGBColor(0x00, 0xd4, 0xff)   # 시안 (메인)
GREEN    = RGBColor(0x00, 0xff, 0x88)   # 그린
PURPLE   = RGBColor(0xcc, 0x44, 0xff)   # 퍼플
YELLOW   = RGBColor(0xff, 0xd7, 0x00)   # 옐로
RED      = RGBColor(0xff, 0x44, 0x66)   # 레드
WHITE    = RGBColor(0xff, 0xff, 0xff)   # 화이트
GRAY     = RGBColor(0x77, 0x77, 0x99)   # 회색
DARK_GRAY= RGBColor(0x33, 0x33, 0x55)   # 진한 회색

W = Inches(13.33)   # 와이드 슬라이드 가로
H = Inches(7.5)     # 와이드 슬라이드 세로

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # 완전 빈 레이아웃


# ════════════════════════════════════════════════════════════════
# 헬퍼 함수
# ════════════════════════════════════════════════════════════════

def add_slide():
    return prs.slides.add_slide(blank)

def bg(slide, color=BG):
    """슬라이드 배경 채우기"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, radius=False):
    """사각형 추가"""
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # 테두리 제거
    return shape

def border_rect(slide, x, y, w, h, border_color, fill_color=None, border_width=Pt(2)):
    """테두리 사각형"""
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    return shape

def txbox(slide, text, x, y, w, h,
          font_size=Pt(18), color=WHITE, bold=False, align=PP_ALIGN.LEFT,
          font_name="맑은 고딕", wrap=True):
    """텍스트 박스 추가"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return tb

def card_box(slide, x, y, w, h, title, desc, icon,
             title_color=CYAN, border_color=CYAN):
    """기능 카드 박스"""
    # 배경
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)

    # 아이콘
    txbox(slide, icon, x + Cm(0.5), y + Cm(0.4), Cm(2), Cm(1.5),
          font_size=Pt(26), color=WHITE, align=PP_ALIGN.CENTER)

    # 제목
    txbox(slide, title, x + Cm(0.5), y + Cm(1.8), w - Cm(1), Cm(0.9),
          font_size=Pt(14), color=title_color, bold=True)

    # 설명
    txbox(slide, desc, x + Cm(0.5), y + Cm(2.6), w - Cm(1), h - Cm(3),
          font_size=Pt(10), color=GRAY)

def divider(slide, x, y, w, color):
    """구분선"""
    shape = slide.shapes.add_shape(1, x, y, w, Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ════════════════════════════════════════════════════════════════
# 슬라이드 1 — 표지
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

# 배경 장식 (글로우 효과 대신 반투명 레이어)
r = rect(s, Inches(0), Inches(0), W, H, BG)

# 상단 데코 바
rect(s, Inches(0), Inches(0), W, Inches(0.12), CYAN)

# 중앙 카드
card_bg = rect(s, Inches(2.5), Inches(1.2), Inches(8.3), Inches(5.1), CARD)
card_bg.line.color.rgb = CYAN
card_bg.line.width = Pt(2)

# 메인 이모지
txbox(s, "🎮", Inches(5.8), Inches(1.5), Inches(1.7), Inches(1.5),
      font_size=Pt(72), align=PP_ALIGN.CENTER)

# 메인 타이틀
txbox(s, "가라온 브로즈", Inches(2.8), Inches(2.8), Inches(7.7), Inches(1.0),
      font_size=Pt(40), color=CYAN, bold=True, align=PP_ALIGN.CENTER,
      font_name="맑은 고딕")

# 서브타이틀
txbox(s, "보드게임카페 전용 스마트 앱", Inches(2.8), Inches(3.7), Inches(7.7), Inches(0.7),
      font_size=Pt(20), color=YELLOW, bold=True, align=PP_ALIGN.CENTER)

# 설명
txbox(s, "게임 추천 · 방문자 관리 · 음악 플레이어 · 집주인 호출 · 게임 도구 8종",
      Inches(2.8), Inches(4.3), Inches(7.7), Inches(0.7),
      font_size=Pt(13), color=GRAY, align=PP_ALIGN.CENTER)

# 태그들
tags = [("🟣 Supabase", PURPLE), ("🟢 Netlify", GREEN), ("🔵 n8n", CYAN), ("🟡 Telegram", YELLOW)]
for i, (tag, col) in enumerate(tags):
    txbox(s, tag, Inches(3.2 + i * 1.8), Inches(5.0), Inches(1.6), Inches(0.5),
          font_size=Pt(11), color=col, align=PP_ALIGN.CENTER)

# 하단 바
rect(s, Inches(0), H - Inches(0.12), W, Inches(0.12), CYAN)

txbox(s, "2026. 03  ·  by JK Networks", Inches(0), H - Inches(0.5), W, Inches(0.4),
      font_size=Pt(10), color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# 슬라이드 2 — 목차
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.12), CYAN)

txbox(s, "📋  목차", Inches(0.6), Inches(0.3), W, Inches(0.7),
      font_size=Pt(28), color=CYAN, bold=True)
divider(s, Inches(0.6), Inches(1.0), Inches(12.1), DARK_GRAY)

items = [
    ("01", "프로젝트 소개",      "무엇을 왜 만들었는가",          CYAN),
    ("02", "서비스 구성",        "6개 화면 & 기술 스택",          GREEN),
    ("03", "게임 추천 시스템",   "인원수·연령대 기반 필터링",      PURPLE),
    ("04", "게임 DB 관리 (GBHQ)","보드게임·맨손게임·호출메뉴",    YELLOW),
    ("05", "방문자 관리",        "QR 체크인·방문기록·단체사진",    RED),
    ("06", "게임 도구 8종",      "사다리·퀴즈·타이머 등",          CYAN),
    ("07", "음악 플레이어",      "Spotify 연동·페이지 이동 재생유지",GREEN),
    ("08", "집주인 호출 & 알림", "텔레그램 실시간 연동",           PURPLE),
]

cols = [Inches(0.7), Inches(7.0)]
for i, (num, title, desc, col) in enumerate(items):
    row = i % 4
    cx = cols[i // 4]
    cy = Inches(1.2) + row * Inches(1.4)

    txbox(s, num, cx, cy + Inches(0.05), Inches(0.7), Inches(0.55),
          font_size=Pt(20), color=col, bold=True, align=PP_ALIGN.CENTER)
    txbox(s, title, cx + Inches(0.75), cy, Inches(4.0), Inches(0.55),
          font_size=Pt(15), color=WHITE, bold=True)
    txbox(s, desc, cx + Inches(0.75), cy + Inches(0.5), Inches(5.5), Inches(0.5),
          font_size=Pt(11), color=GRAY)
    divider(s, cx, cy + Inches(1.2), Inches(5.8), DARK_GRAY)


# ════════════════════════════════════════════════════════════════
# 슬라이드 3 — 프로젝트 소개
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.12), GREEN)

txbox(s, "01  |  프로젝트 소개", Inches(0.6), Inches(0.3), W, Inches(0.6),
      font_size=Pt(24), color=GREEN, bold=True)
divider(s, Inches(0.6), Inches(0.95), Inches(12.1), DARK_GRAY)

# 왼쪽: 배경/목적
txbox(s, "🏪  개발 배경", Inches(0.6), Inches(1.1), Inches(5.5), Inches(0.55),
      font_size=Pt(16), color=CYAN, bold=True)
txbox(s,
      "가라온·라온 형제를 위한 프라이빗 보드게임카페를 운영하면서\n"
      "게임 추천, 방문자 관리, 주문 호출 등을 수작업으로 처리하는\n"
      "비효율을 해소하고자 전용 앱을 개발하였습니다.",
      Inches(0.6), Inches(1.65), Inches(5.8), Inches(1.5),
      font_size=Pt(13), color=WHITE)

txbox(s, "🎯  핵심 목표", Inches(0.6), Inches(3.2), Inches(5.5), Inches(0.55),
      font_size=Pt(16), color=CYAN, bold=True)
goals = [
    "✅  인원수·나이에 맞는 게임 자동 추천",
    "✅  방문자 체크인/아웃 및 방문기록 카드 관리",
    "✅  아이들 ↔ 집주인 실시간 텔레그램 호출",
    "✅  사다리·퀴즈·타이머 등 게임 보조 도구 제공",
    "✅  배경음악 Spotify 재생 (페이지 이동 중에도 유지)",
]
for j, g in enumerate(goals):
    txbox(s, g, Inches(0.6), Inches(3.75) + j * Inches(0.52), Inches(5.8), Inches(0.5),
          font_size=Pt(12), color=WHITE)

# 오른쪽: 스펙 카드
border_rect(s, Inches(7.0), Inches(1.1), Inches(5.7), Inches(5.8), CYAN, CARD)
txbox(s, "⚙️  기술 스택", Inches(7.2), Inches(1.2), Inches(5.3), Inches(0.55),
      font_size=Pt(15), color=CYAN, bold=True)

specs = [
    ("프론트엔드",  "순수 HTML · CSS · JavaScript (No Framework)"),
    ("데이터베이스","Supabase (PostgreSQL + Storage)"),
    ("음악",        "Spotify Embed API  +  iframe SPA 방식"),
    ("알림",        "Telegram Bot API  (양방향 메시지)"),
    ("배포",        "Netlify  →  office-ai.app/garaon-bros/"),
    ("자동화",      "n8n Cloud  (Webhook 연동)"),
    ("이미지",      "Wikipedia API  +  QR 촬영 업로드"),
    ("AI 검색",     "게임 정보 자동 입력 (ChatGPT 연동)"),
]
for j, (k, v) in enumerate(specs):
    cy = Inches(1.85) + j * Inches(0.58)
    txbox(s, k, Inches(7.2), cy, Inches(1.8), Inches(0.5),
          font_size=Pt(11), color=YELLOW, bold=True)
    txbox(s, v, Inches(9.0), cy, Inches(3.5), Inches(0.5),
          font_size=Pt(11), color=WHITE)


# ════════════════════════════════════════════════════════════════
# 슬라이드 4 — 서비스 구성 (화면 맵)
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.12), PURPLE)

txbox(s, "02  |  서비스 구성", Inches(0.6), Inches(0.3), W, Inches(0.6),
      font_size=Pt(24), color=PURPLE, bold=True)
divider(s, Inches(0.6), Inches(0.95), Inches(12.1), DARK_GRAY)

screens = [
    ("🏠", "index.html",      "메인 허브",       "체크인·메뉴 진입점",      CYAN),
    ("🎯", "recommend.html",  "게임 추천",       "인원수·연령대 필터",      GREEN),
    ("📊", "gbhq.html",       "게임 DB 관리",    "게임 등록·편집·호출메뉴", YELLOW),
    ("🛠️", "game-tools.html", "게임 도구 8종",   "사다리·퀴즈·타이머 등",  PURPLE),
    ("📝", "review.html",     "게임 후기",       "별점·코멘트 기록",        RED),
    ("📸", "upload.html",     "QR 촬영 업로드",  "방문자 프로필 사진",      CYAN),
]

for i, (icon, fname, title, desc, col) in enumerate(screens):
    col_idx = i % 3
    row_idx = i // 3
    cx = Inches(0.5) + col_idx * Inches(4.3)
    cy = Inches(1.1) + row_idx * Inches(2.8)
    cw = Inches(4.0)
    ch = Inches(2.5)

    bg_shape = slide_shape = s.shapes.add_shape(1, cx, cy, cw, ch)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = CARD
    bg_shape.line.color.rgb = col
    bg_shape.line.width = Pt(2)

    txbox(s, icon, cx + Inches(0.2), cy + Inches(0.15), Inches(0.9), Inches(0.8),
          font_size=Pt(28), align=PP_ALIGN.CENTER)
    txbox(s, title, cx + Inches(1.1), cy + Inches(0.15), Inches(2.7), Inches(0.55),
          font_size=Pt(14), color=col, bold=True)
    txbox(s, fname, cx + Inches(1.1), cy + Inches(0.65), Inches(2.7), Inches(0.4),
          font_size=Pt(9), color=DARK_GRAY)
    txbox(s, desc, cx + Inches(0.2), cy + Inches(1.2), Inches(3.7), Inches(0.9),
          font_size=Pt(11), color=WHITE)

# 오른쪽 중앙: 공유 모듈
txbox(s, "🔧  공유 모듈", Inches(12.7), Inches(1.5), Inches(0.5), Inches(2.0),
      font_size=Pt(10), color=GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# 슬라이드 5 — 게임 추천 시스템
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.12), PURPLE)

txbox(s, "03  |  게임 추천 시스템", Inches(0.6), Inches(0.3), W, Inches(0.6),
      font_size=Pt(24), color=PURPLE, bold=True)
divider(s, Inches(0.6), Inches(0.95), Inches(12.1), DARK_GRAY)

# 왼쪽 설명
txbox(s, "🎯  추천 흐름", Inches(0.6), Inches(1.1), Inches(5.5), Inches(0.55),
      font_size=Pt(16), color=PURPLE, bold=True)

flow = [
    ("1️⃣", "인원수 선택", "2명 ~ 10명+ 버튼 탭"),
    ("2️⃣", "연령대 필터", "아동(~10) / 청소년(11~17) / 성인(18+) / 전체"),
    ("3️⃣", "카테고리",   "전략·카드·파티·협력·추리·경제·공포·롤플레이 등"),
    ("4️⃣", "결과 카드",  "이미지·난이도·플레이시간·재미도 표시"),
    ("5️⃣", "상세 보기",  "규칙·팁·YouTube 영상·AI 검색 자동 입력"),
]

for j, (num, title, desc) in enumerate(flow):
    cy = Inches(1.7) + j * Inches(0.95)
    border_rect(s, Inches(0.6), cy, Inches(5.8), Inches(0.82), PURPLE, CARD)
    txbox(s, num, Inches(0.7), cy + Inches(0.12), Inches(0.6), Inches(0.6),
          font_size=Pt(18), align=PP_ALIGN.CENTER)
    txbox(s, title, Inches(1.4), cy + Inches(0.05), Inches(2.5), Inches(0.45),
          font_size=Pt(13), color=WHITE, bold=True)
    txbox(s, desc, Inches(1.4), cy + Inches(0.44), Inches(4.8), Inches(0.35),
          font_size=Pt(10), color=GRAY)

# 오른쪽: DB 구조
border_rect(s, Inches(7.0), Inches(1.1), Inches(5.7), Inches(5.8), PURPLE, CARD)
txbox(s, "🗄️  DB 테이블 구조", Inches(7.2), Inches(1.2), Inches(5.3), Inches(0.55),
      font_size=Pt(15), color=PURPLE, bold=True)

fields_bg = [
    ("board_games", [
        "name / name_en",
        "min_players / max_players",
        "play_time_min / difficulty / fun_score",
        "categories[]  age_min",
        "image_url / description",
        "rules / tips / youtube_url",
    ]),
    ("hand_games", [
        "name / min_players / max_players",
        "categories[]  age_min",
        "description / rules / tips",
        "materials (도구 없이 하는 게임)",
    ]),
]

cy = Inches(1.85)
for tbl, fields in fields_bg:
    txbox(s, f"📦 {tbl}", Inches(7.2), cy, Inches(5.3), Inches(0.45),
          font_size=Pt(12), color=YELLOW, bold=True)
    cy += Inches(0.42)
    for f in fields:
        txbox(s, f"  ·  {f}", Inches(7.2), cy, Inches(5.3), Inches(0.38),
              font_size=Pt(10), color=WHITE)
        cy += Inches(0.36)
    cy += Inches(0.15)


# ════════════════════════════════════════════════════════════════
# 슬라이드 6 — 게임 DB 관리 (GBHQ)
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.12), YELLOW)

txbox(s, "04  |  게임 DB 관리  (GBHQ)", Inches(0.6), Inches(0.3), W, Inches(0.6),
      font_size=Pt(24), color=YELLOW, bold=True)
divider(s, Inches(0.6), Inches(0.95), Inches(12.1), DARK_GRAY)

features = [
    ("🎲", "보드게임 관리",
     "· 게임 목록 카드형 UI\n· 수정 아이콘(우상단) + 상세보기\n· AI 검색으로 이름 입력 시 자동 채우기\n· Wikipedia 이미지 자동 가져오기\n· YouTube 인라인 재생",
     CYAN),
    ("✋", "맨손 게임",
     "· 도구 없이 사람끼리 하는 게임\n· 별도 테이블 (hand_games)\n· 카테고리: 언어·신체·추리·창의\n· 인원수·연령대 필터 동일 적용",
     GREEN),
    ("🧑‍🤝‍🧑", "방문자·방문기록",
     "· QR 촬영으로 프로필 사진 업로드\n· 이모지 아바타 선택\n· 방문기록 카드형 UI (3개월~전체)\n· 단체사진 업로드 (최대 5장)\n· 방문기록 ↔ 게임 호출 이력 연동",
     PURPLE),
    ("📋", "호출 메뉴 관리",
     "· 메뉴 아이템 등록 (이름·이모지)\n· 솔드아웃(품절) 표시 기능\n· 아이들 → 텔레그램으로 집주인 호출\n· 집주인 → 아이들 메시지 수신\n· 확인 시 자동 텔레그램 답장",
     YELLOW),
]

for i, (icon, title, desc, col) in enumerate(features):
    cx = Inches(0.5) + (i % 2) * Inches(6.4)
    cy = Inches(1.1) + (i // 2) * Inches(2.9)
    cw = Inches(6.0)
    ch = Inches(2.6)
    border_rect(s, cx, cy, cw, ch, col, CARD)
    txbox(s, icon, cx + Inches(0.2), cy + Inches(0.1), Inches(0.8), Inches(0.7),
          font_size=Pt(28))
    txbox(s, title, cx + Inches(1.1), cy + Inches(0.1), Inches(4.7), Inches(0.55),
          font_size=Pt(14), color=col, bold=True)
    txbox(s, desc, cx + Inches(0.25), cy + Inches(0.75), Inches(5.5), Inches(1.7),
          font_size=Pt(11), color=WHITE)


# ════════════════════════════════════════════════════════════════
# 슬라이드 7 — 방문자 관리
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.12), RED)

txbox(s, "05  |  방문자 관리", Inches(0.6), Inches(0.3), W, Inches(0.6),
      font_size=Pt(24), color=RED, bold=True)
divider(s, Inches(0.6), Inches(0.95), Inches(12.1), DARK_GRAY)

# 체크인 흐름 (왼쪽)
txbox(s, "📲  체크인 흐름", Inches(0.6), Inches(1.1), Inches(5.5), Inches(0.55),
      font_size=Pt(16), color=RED, bold=True)

checkin_steps = [
    ("이름 입력", "기존 방문자 자동완성 지원"),
    ("나이 입력", "연령대 필터에 자동 반영"),
    ("사진 방식 선택", "① QR 촬영  ② 이모지 아바타 선택"),
    ("QR 촬영", "일회용 토큰 생성 → 폰 카메라 → Supabase Storage 업로드"),
    ("체크인 완료", "방문기록 자동 생성 + 환영 팝업"),
]
for j, (title, desc) in enumerate(checkin_steps):
    cy = Inches(1.7) + j * Inches(0.88)
    txbox(s, f"{j+1}.", Inches(0.6), cy, Inches(0.4), Inches(0.5),
          font_size=Pt(14), color=RED, bold=True)
    txbox(s, title, Inches(1.05), cy, Inches(2.5), Inches(0.42),
          font_size=Pt(13), color=WHITE, bold=True)
    txbox(s, desc, Inches(1.05), cy + Inches(0.4), Inches(4.9), Inches(0.4),
          font_size=Pt(10), color=GRAY)
    divider(s, Inches(0.6), cy + Inches(0.8), Inches(5.8), DARK_GRAY)

# 오른쪽: 방문기록 기능
border_rect(s, Inches(7.0), Inches(1.1), Inches(5.7), Inches(5.8), RED, CARD)
txbox(s, "📋  방문기록 기능", Inches(7.2), Inches(1.2), Inches(5.3), Inches(0.55),
      font_size=Pt(15), color=RED, bold=True)

visit_features = [
    ("카드형 UI", "방문자 사진·이름·호출이력 카드로 표시"),
    ("기간 필터", "3개월 / 6개월 / 1년 / 전체 선택"),
    ("호출 이력", "item_label·item_emoji·note 컬럼으로 기록"),
    ("단체사진", "QR 촬영 방식, 최대 5장 업로드"),
    ("사진 삭제", "단체사진 개별 삭제 버튼(X)"),
    ("체크아웃", "하단 바에서 선택 인원 일괄 체크아웃"),
]
for j, (title, desc) in enumerate(visit_features):
    cy = Inches(1.85) + j * Inches(0.8)
    txbox(s, f"·  {title}", Inches(7.2), cy, Inches(2.2), Inches(0.42),
          font_size=Pt(12), color=CYAN, bold=True)
    txbox(s, desc, Inches(9.4), cy, Inches(3.0), Inches(0.42),
          font_size=Pt(11), color=WHITE)
    divider(s, Inches(7.2), cy + Inches(0.72), Inches(5.3), DARK_GRAY)


# ════════════════════════════════════════════════════════════════
# 슬라이드 8 — 게임 도구 8종
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.12), CYAN)

txbox(s, "06  |  게임 도구  8종", Inches(0.6), Inches(0.3), W, Inches(0.6),
      font_size=Pt(24), color=CYAN, bold=True)
divider(s, Inches(0.6), Inches(0.95), Inches(12.1), DARK_GRAY)

tools = [
    ("🪜", "사다리 게임",    "참여자 이름 입력 → 자동 사다리 생성\n결과 애니메이션으로 공정하게 결정"),
    ("🧩", "보드게임 퀴즈",  "카테고리별 퀴즈 + 정지 버튼\n타이머 연동으로 긴장감 연출"),
    ("⏱️", "게임 타이머",    "카운트다운 + 경고음\n다음 라운드 자동 리셋"),
    ("🎲", "주사위",         "1~6 랜덤 굴리기\n다중 주사위 동시 지원"),
    ("🃏", "카드 셔플",      "카드 덱 가상 셔플\n원하는 장수 뽑기"),
    ("🔢", "점수판",         "팀별·인원별 점수 기록\n라운드별 히스토리"),
    ("🎵", "BGM 플레이어",   "Spotify 통합\n페이지 이동 중에도 재생 유지"),
    ("🔍", "게임 검색",      "실시간 필터링\nDB 전체 게임 빠른 탐색"),
]

for i, (icon, title, desc) in enumerate(tools):
    col_idx = i % 4
    row_idx = i // 4
    cx = Inches(0.4) + col_idx * Inches(3.2)
    cy = Inches(1.1) + row_idx * Inches(2.8)
    cw = Inches(3.0)
    ch = Inches(2.5)
    col = [CYAN, GREEN, YELLOW, PURPLE][i % 4]

    border_rect(s, cx, cy, cw, ch, col, CARD)
    txbox(s, icon, cx + Inches(0.15), cy + Inches(0.1), Inches(0.9), Inches(0.8),
          font_size=Pt(28), align=PP_ALIGN.CENTER)
    txbox(s, f"{i+1}. {title}", cx + Inches(0.15), cy + Inches(0.85), Inches(2.7), Inches(0.5),
          font_size=Pt(12), color=col, bold=True)
    txbox(s, desc, cx + Inches(0.15), cy + Inches(1.35), Inches(2.7), Inches(1.0),
          font_size=Pt(10), color=WHITE)


# ════════════════════════════════════════════════════════════════
# 슬라이드 9 — 음악 플레이어
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.12), GREEN)

txbox(s, "07  |  음악 플레이어  (Spotify 연동)", Inches(0.6), Inches(0.3), W, Inches(0.6),
      font_size=Pt(24), color=GREEN, bold=True)
divider(s, Inches(0.6), Inches(0.95), Inches(12.1), DARK_GRAY)

# 좌측 설명
txbox(s, "🎵  핵심 기술: iframe SPA 방식", Inches(0.6), Inches(1.1), Inches(6.0), Inches(0.55),
      font_size=Pt(16), color=GREEN, bold=True)
txbox(s,
      "일반적인 SPA에서 페이지를 이동하면 Spotify 재생이 끊깁니다.\n"
      "이를 해결하기 위해 메인 페이지를 Spotify 플레이어 호스트로 고정하고,\n"
      "나머지 콘텐츠를 <iframe> 내에서 로드하는 방식을 적용했습니다.",
      Inches(0.6), Inches(1.75), Inches(6.0), Inches(1.3),
      font_size=Pt(12), color=WHITE)

music_features = [
    ("🔍", "검색", "곡·앨범·플레이리스트 통합 검색"),
    ("▶️", "재생", "Spotify Embed API 자동 연동"),
    ("📱", "반응형", "음악 배너 동적 레이아웃 (빔프로젝터 모드)"),
    ("🔗", "유지",  "iframe 네비게이션으로 음악 끊김 없이 페이지 이동"),
    ("🎶", "로그인", "비회원도 미리듣기 / 회원은 전곡 재생"),
]
for j, (icon, title, desc) in enumerate(music_features):
    cy = Inches(3.2) + j * Inches(0.72)
    border_rect(s, Inches(0.6), cy, Inches(6.0), Inches(0.62), GREEN, CARD)
    txbox(s, icon, Inches(0.7), cy + Inches(0.08), Inches(0.5), Inches(0.45),
          font_size=Pt(16), align=PP_ALIGN.CENTER)
    txbox(s, title, Inches(1.25), cy + Inches(0.08), Inches(1.4), Inches(0.45),
          font_size=Pt(12), color=GREEN, bold=True)
    txbox(s, desc, Inches(2.7), cy + Inches(0.08), Inches(3.8), Inches(0.45),
          font_size=Pt(11), color=WHITE)

# 우측: 코드 흐름 다이어그램
border_rect(s, Inches(7.0), Inches(1.1), Inches(5.7), Inches(5.8), GREEN, CARD)
txbox(s, "🔄  페이지 이동 흐름", Inches(7.2), Inches(1.2), Inches(5.3), Inches(0.55),
      font_size=Pt(15), color=GREEN, bold=True)

flow_items = [
    ("메인 페이지 로드", "Spotify iframe 초기화 + 음악 재생 시작", GREEN),
    ("콘텐츠 영역",      "별도 <iframe>으로 서브 페이지 렌더링",    CYAN),
    ("링크 클릭",        "location.href 인터셉트 → iframe.src 변경", YELLOW),
    ("서브 페이지 이동", "메인 Spotify 플레이어 영향 없음",          GREEN),
    ("카드 클릭 방식",   "onclick 직접 호출로 음악 끊김 완전 방지",  PURPLE),
]
for j, (title, desc, col) in enumerate(flow_items):
    cy = Inches(1.85) + j * Inches(0.95)
    txbox(s, f"{'↓' if j > 0 else '▶'} {title}", Inches(7.2), cy, Inches(5.3), Inches(0.42),
          font_size=Pt(12), color=col, bold=True)
    txbox(s, desc, Inches(7.5), cy + Inches(0.4), Inches(4.9), Inches(0.38),
          font_size=Pt(10), color=GRAY)


# ════════════════════════════════════════════════════════════════
# 슬라이드 10 — 집주인 호출 & 텔레그램
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.12), PURPLE)

txbox(s, "08  |  집주인 호출  &  텔레그램 연동", Inches(0.6), Inches(0.3), W, Inches(0.6),
      font_size=Pt(24), color=PURPLE, bold=True)
divider(s, Inches(0.6), Inches(0.95), Inches(12.1), DARK_GRAY)

# 양방향 메시지 다이어그램
txbox(s, "🔄  양방향 실시간 메시지", Inches(0.6), Inches(1.1), Inches(6.0), Inches(0.55),
      font_size=Pt(16), color=PURPLE, bold=True)

# 아이들 → 집주인
border_rect(s, Inches(0.6), Inches(1.75), Inches(5.8), Inches(2.3), PURPLE, CARD)
txbox(s, "📱 아이들  →  집주인", Inches(0.8), Inches(1.85), Inches(5.4), Inches(0.5),
      font_size=Pt(14), color=PURPLE, bold=True)
txbox(s,
      "① 호출 메뉴에서 아이템 탭 (ex: 🍕 피자 주세요)\n"
      "② 앱에서 Telegram Bot API 직접 호출\n"
      "③ 집주인 텔레그램으로 주문 내용 도착\n"
      "④ 방문기록에 호출 이력 자동 저장",
      Inches(0.8), Inches(2.35), Inches(5.4), Inches(1.6),
      font_size=Pt(11), color=WHITE)

# 집주인 → 아이들
border_rect(s, Inches(0.6), Inches(4.2), Inches(5.8), Inches(2.1), CYAN, CARD)
txbox(s, "💬 집주인  →  아이들", Inches(0.8), Inches(4.3), Inches(5.4), Inches(0.5),
      font_size=Pt(14), color=CYAN, bold=True)
txbox(s,
      "① 집주인이 텔레그램에서 메시지 전송\n"
      "② 앱이 n8n Webhook으로 메시지 수신\n"
      "③ 모달 팝업 + 알림음 반복 재생\n"
      "④ 확인 버튼 → 텔레그램 자동 답장",
      Inches(0.8), Inches(4.8), Inches(5.4), Inches(1.4),
      font_size=Pt(11), color=WHITE)

# 우측: 솔드아웃 & 추가 기능
border_rect(s, Inches(7.0), Inches(1.1), Inches(5.7), Inches(5.8), YELLOW, CARD)
txbox(s, "🛒  호출 메뉴 특별 기능", Inches(7.2), Inches(1.2), Inches(5.3), Inches(0.55),
      font_size=Pt(15), color=YELLOW, bold=True)

call_features = [
    ("솔드아웃(품절)", "관리자가 품절 처리 → 아이들 화면에\n빨간 SOLD OUT 뱃지 표시, 탭 비활성화", YELLOW),
    ("이모지 피커",   "메뉴 등록 시 이모지 선택 팝업\n시각적으로 구분하기 쉽게", GREEN),
    ("실시간 알림음", "집주인 메시지 수신 시 알림음 반복\n확인 전까지 계속 재생", RED),
    ("메시지 히스토리","호출 내역을 방문기록과 연동\n item_label·item_emoji·note 저장", CYAN),
]
for j, (title, desc, col) in enumerate(call_features):
    cy = Inches(1.85) + j * Inches(1.2)
    txbox(s, f"·  {title}", Inches(7.2), cy, Inches(5.3), Inches(0.45),
          font_size=Pt(13), color=col, bold=True)
    txbox(s, desc, Inches(7.5), cy + Inches(0.43), Inches(5.0), Inches(0.65),
          font_size=Pt(10), color=WHITE)
    divider(s, Inches(7.2), cy + Inches(1.1), Inches(5.3), DARK_GRAY)


# ════════════════════════════════════════════════════════════════
# 슬라이드 11 — 마무리 & URL
# ════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
rect(s, Inches(0), Inches(0), W, Inches(0.12), CYAN)

# 배경 카드
border_rect(s, Inches(1.5), Inches(0.8), Inches(10.3), Inches(5.9), CYAN, CARD)

txbox(s, "🎮", Inches(5.9), Inches(0.9), Inches(1.5), Inches(1.3),
      font_size=Pt(60), align=PP_ALIGN.CENTER)

txbox(s, "가라온 브로즈", Inches(1.8), Inches(2.0), Inches(9.7), Inches(0.9),
      font_size=Pt(36), color=CYAN, bold=True, align=PP_ALIGN.CENTER)

txbox(s, "보드게임카페 전용 스마트 앱  —  완성",
      Inches(1.8), Inches(2.85), Inches(9.7), Inches(0.6),
      font_size=Pt(16), color=YELLOW, align=PP_ALIGN.CENTER)

# URL
border_rect(s, Inches(2.5), Inches(3.6), Inches(8.3), Inches(0.65), GREEN, CARD2)
txbox(s, "🌐  https://office-ai.app/garaon-bros/",
      Inches(2.6), Inches(3.65), Inches(8.1), Inches(0.55),
      font_size=Pt(14), color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# 향후 계획
txbox(s, "💡  향후 계획", Inches(2.0), Inches(4.45), Inches(9.3), Inches(0.5),
      font_size=Pt(14), color=GRAY, bold=True)
plans = [
    "· 게임 리뷰 별점 누적 통계 대시보드",
    "· 맞춤형 게임 추천 AI 강화 (플레이 이력 기반)",
    "· 방문자 포인트/도장 적립 시스템",
    "· PWA 오프라인 지원",
]
for j, p in enumerate(plans):
    txbox(s, p, Inches(2.2), Inches(4.95) + j * Inches(0.42), Inches(9.0), Inches(0.4),
          font_size=Pt(11), color=GRAY)

rect(s, Inches(0), H - Inches(0.12), W, Inches(0.12), CYAN)
txbox(s, "JK Networks  ·  2026.03", Inches(0), H - Inches(0.45), W, Inches(0.35),
      font_size=Pt(10), color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════════
output_path = "C:/work/office-ai/garaon-bros/가라온브로즈-소개자료.pptx"
prs.save(output_path)
print(f"✅  저장 완료: {output_path}")
print(f"    슬라이드 수: {len(prs.slides)}장")
